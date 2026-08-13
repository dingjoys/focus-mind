#!/usr/bin/env python3
"""Build the LaunchMinds 'From Assistant to Autonomous' presentation (PPTX)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- Brand palette (dark, indigo/mint accents) ----
BG_DARK = RGBColor(0x0B, 0x0F, 0x1A)
BG_PANEL = RGBColor(0x14, 0x1B, 0x2E)
BG_PANEL2 = RGBColor(0x18, 0x21, 0x38)
ACCENT = RGBColor(0x7C, 0x9C, 0xFF)      # indigo
ACCENT2 = RGBColor(0x5E, 0xE6, 0xC0)     # mint
TEXT_MAIN = RGBColor(0xF2, 0xF4, 0xF8)
TEXT_DIM = RGBColor(0xA6, 0xAF, 0xC2)
TEXT_FAINT = RGBColor(0x6B, 0x74, 0x8C)

FONT = "Avenir Next"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_bg(slide, color=BG_DARK):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color, line_color=None, radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, x, y, w, h)
    if radius:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=18, color=TEXT_MAIN, bold=False,
             align=PP_ALIGN.LEFT, font=FONT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.15, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font
    return tb


def add_bullets(slide, x, y, w, h, items, size=15, color=TEXT_MAIN,
                 font=FONT, space_after=10, line_spacing=1.2):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = f"—  {item}"
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = font
    return tb


def add_kicker(slide, text, num_text=None):
    add_rect(slide, Inches(0.55), Inches(0.52), Inches(0.42), Inches(0.055), ACCENT2)
    add_text(slide, Inches(0.55), Inches(0.62), Inches(8), Inches(0.4),
             text.upper(), size=12.5, color=TEXT_DIM, bold=True)
    if num_text:
        add_text(slide, Inches(11.9), Inches(0.55), Inches(0.9), Inches(0.4),
                 num_text, size=12.5, color=TEXT_FAINT, align=PP_ALIGN.RIGHT)


def add_footer(slide, text="LaunchMinds  ·  Built on Minds  ·  2026"):
    add_text(slide, Inches(0.55), Inches(7.1), Inches(9), Inches(0.3),
             text, size=9.5, color=TEXT_FAINT)


def add_badge(slide, x, y, text, color):
    w = Inches(0.16 * len(text) + 0.5)
    shp = add_rect(slide, x, y, w, Inches(0.32), BG_DARK, line_color=color, radius=0.5)
    tf = shp.text_frame
    tf.word_wrap = False
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = color
    r.font.name = FONT
    return shp


def new_slide(bg=BG_DARK):
    s = prs.slides.add_slide(BLANK)
    add_bg(s, bg)
    return s


def set_notes(slide, text):
    notes = slide.notes_slide
    notes.notes_text_frame.text = text


# ---------------------------------------------------------------------------
# Slide 1 — Title
# ---------------------------------------------------------------------------
s = new_slide()
add_rect(s, Inches(0), Inches(0), Inches(0.14), SLIDE_H, ACCENT2)
add_text(s, Inches(0.9), Inches(1.85), Inches(11), Inches(0.5),
         "LAUNCHMINDS", size=20, color=ACCENT2, bold=True)
add_text(s, Inches(0.85), Inches(2.3), Inches(11.5), Inches(1.9),
         "From Assistant\nto Autonomous", size=54, color=TEXT_MAIN, bold=True, line_spacing=1.05)
add_text(s, Inches(0.9), Inches(4.3), Inches(10.5), Inches(0.6),
         "Built on Minds, from day one.",
         size=19, color=TEXT_DIM, italic=True)
add_bullets(s, Inches(0.9), Inches(5.1), Inches(10.8), Inches(1.8), [
    "Agentic Campaign Operations System for project teams",
    "Built on Minds' persistent agent platform",
    "Plan → Approve → Execute → Verify → Settle → Learn",
], size=15, color=TEXT_DIM, space_after=8)
add_footer(s, "LaunchMinds · 2026")
set_notes(s, "Good afternoon. LaunchMinds is the Agentic Campaign Operations System for project teams. "
             "Instead of helping with one isolated marketing task, it maintains project intelligence, carries "
             "campaign work across sessions, and moves from an objective to a verified outcome within explicit "
             "approval boundaries. We built it on Minds because campaign operations are long-running, stateful, "
             "multi-party, and increasingly transactional. Our name reflects that directly — Launch, plus "
             "Minds — because everything we're about to show you is built on your platform, not next to it. "
             "Today I want to show you how we're moving from assisted work to persistent, accountable campaign "
             "operations.")

# ---------------------------------------------------------------------------
# Slide 2 — The Problem
# ---------------------------------------------------------------------------
s = new_slide()
add_kicker(s, "The Problem", "01 / 07")
add_text(s, Inches(0.55), Inches(1.05), Inches(11.5), Inches(1.3),
         "Campaign operations today are\nstitched together by hand", size=34, color=TEXT_MAIN, bold=True, line_spacing=1.1)
add_bullets(s, Inches(0.55), Inches(2.85), Inches(11.5), Inches(3.2), [
    "One tool for campaign design, another for deployment, a spreadsheet for the rest",
    "Campaign design, deployment, resource orchestration — all disconnected",
    "Nothing persists between sessions, and nothing is accountable across the lifecycle",
    "The fix isn't another point tool — it's persistent, accountable campaign operations",
], size=19, color=TEXT_MAIN, space_after=18)
add_footer(s)
set_notes(s, "Campaign operations today are stitched together by hand — one tool for campaign design, another "
             "for deployment, a spreadsheet for resource orchestration, and a person holding it all in their "
             "head. Nothing persists between sessions, and nothing is accountable across the full lifecycle, "
             "from objective to verified outcome. We don't think the fix is another point tool. It's a system "
             "that carries state, enforces approval, and stays accountable across the entire campaign "
             "lifecycle.")

# ---------------------------------------------------------------------------
# Slide 3 — Positioning: Minds (platform) + LaunchMinds (vertical control plane)
# ---------------------------------------------------------------------------
s = new_slide()
add_kicker(s, "Positioning", "02 / 07")
add_text(s, Inches(0.55), Inches(0.95), Inches(12.2), Inches(0.7),
         "Minds + LaunchMinds: Platform and Vertical Control Plane", size=27, color=TEXT_MAIN, bold=True)

# Minds: horizontal agent platform
minds_y = Inches(1.75)
box_h = Inches(1.15)
minds_box = add_rect(s, Inches(0.55), minds_y, Inches(12.2), box_h, BG_PANEL2, radius=0.08)
add_rect(s, Inches(0.55), minds_y, Inches(0.08), box_h, ACCENT)
add_text(s, Inches(0.9), minds_y + Inches(0.16), Inches(11.3), Inches(0.35),
         "MINDS — horizontal agent platform", size=15, color=ACCENT, bold=True)
add_text(s, Inches(0.9), minds_y + Inches(0.56), Inches(11.3), Inches(0.5),
         "Model routing  ·  identity  ·  memory  ·  Skills  ·  Tools  ·  collaboration  ·  wallet",
         size=14.5, color=TEXT_DIM)

plus_y = minds_y + box_h + Inches(0.06)
add_text(s, Inches(0.55), plus_y, Inches(12.2), Inches(0.3),
         "+", size=20, color=TEXT_FAINT, align=PP_ALIGN.CENTER)

# LaunchMinds: vertical campaign operations control plane
lm_y = plus_y + Inches(0.32)
lm_box = add_rect(s, Inches(0.55), lm_y, Inches(12.2), box_h, BG_PANEL, radius=0.08)
add_rect(s, Inches(0.55), lm_y, Inches(0.08), box_h, ACCENT2)
add_text(s, Inches(0.9), lm_y + Inches(0.16), Inches(11.3), Inches(0.35),
         "LAUNCHMINDS — campaign operations control plane", size=15, color=ACCENT2, bold=True)
add_text(s, Inches(0.9), lm_y + Inches(0.56), Inches(11.3), Inches(0.5),
         "Project intelligence  ·  state  ·  approvals  ·  verification  ·  incentives  ·  settlement  ·  learning",
         size=14.5, color=TEXT_DIM)

tagline_y = lm_y + box_h + Inches(0.28)
add_text(s, Inches(0.55), tagline_y, Inches(12.2), Inches(0.5),
         "Minds provides general agency; LaunchMinds provides campaign accountability.",
         size=17, color=TEXT_MAIN, italic=True)

flow_y = tagline_y + Inches(0.55)
flow_box = add_rect(s, Inches(0.55), flow_y, Inches(12.2), Inches(0.7), BG_PANEL, radius=0.5)
add_text(s, Inches(0.55), flow_y + Inches(0.16), Inches(12.2), Inches(0.4),
         "Objective  →  Approve  →  Execute  →  Verify  →  Settle  →  Learn",
         size=16, color=ACCENT2, bold=True, align=PP_ALIGN.CENTER)
add_footer(s)
set_notes(s, "Minds is not simply the model underneath LaunchMinds. It is the horizontal agent platform that "
             "gives every Mind continuity, identity, memory, tools, collaboration, and the ability to keep "
             "working beyond a single chat. LaunchMinds is the vertical operating layer built on top. We "
             "maintain the trusted project intelligence, campaign state, approval rules, budget constraints, "
             "participant evidence, settlement logic, and outcome history. Minds provides the general ability "
             "to reason and act. LaunchMinds defines how that ability operates safely and measurably inside "
             "campaign operations. That is the division of labor: Minds makes agents persistent and capable; "
             "LaunchMinds makes campaign operations accountable.")

# ---------------------------------------------------------------------------
# Slide 4 — How Agents Evolve (framework intro)
# ---------------------------------------------------------------------------
s = new_slide()
add_kicker(s, "The Framework", "03 / 07")
add_text(s, Inches(0.55), Inches(1.05), Inches(11.5), Inches(0.7),
         "How agents evolve", size=36, color=TEXT_MAIN, bold=True)
add_text(s, Inches(0.55), Inches(1.85), Inches(11.8), Inches(0.9),
         "The same arc coding tools went through.\nWe're climbing it stage by stage — natively on Minds.",
         size=18, color=TEXT_DIM, line_spacing=1.3)

stage_labels = ["Assisted\nCompletion", "Single-Session\nExecution", "Persistent,\nBounded Ops"]
stage_x = [Inches(0.9), Inches(4.95), Inches(9.0)]
box_w = Inches(3.4)
box_y = Inches(3.1)
box_h = Inches(2.5)
colors = [TEXT_FAINT, ACCENT, ACCENT2]
for i, (label, x, c) in enumerate(zip(stage_labels, stage_x, colors)):
    add_rect(s, x, box_y, box_w, box_h, BG_PANEL, radius=0.08)
    add_text(s, x, box_y + Inches(0.3), box_w, Inches(0.4), f"STAGE {i+1}",
             size=13, color=c, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, box_y + Inches(0.85), box_w, Inches(1.2), label,
             size=22, color=TEXT_MAIN, bold=True, align=PP_ALIGN.CENTER, line_spacing=1.15)
    if i < 2:
        add_text(s, x + box_w + Inches(0.05), box_y + Inches(0.95), Inches(0.5), Inches(0.6),
                 "→", size=30, color=TEXT_DIM, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.55), Inches(6.05), Inches(11.8), Inches(0.7),
         "Each stage is a Minds workflow — not a script with a model bolted on.",
         size=17, color=ACCENT2, italic=True)
add_footer(s)
set_notes(s, "If you've watched how coding assistants evolved, you've seen this arc: first, assisted "
             "completion — suggestions, not action. Then single-session execution — an agent that completes a "
             "full task in one sitting, with a human watching. Then persistent, bounded operation — "
             "long-running, accountable, checking its own work within explicit limits. We're building "
             "LaunchMinds through that same three-stage arc, applied to campaign operations instead of code. "
             "And because we built it Minds-native from the start, each stage is a Minds workflow — not a "
             "script with a model bolted on.")


# ---------------------------------------------------------------------------
# Helper for the three side-by-side "stage" slides
# ---------------------------------------------------------------------------
def stage_slide(kicker, num_text, stage_num, stage_name, status_text, status_color,
                 coding_title, coding_bullets, lm_title, lm_bullets, footer_line, notes):
    s = new_slide()
    add_kicker(s, kicker, num_text)
    add_text(s, Inches(0.55), Inches(1.0), Inches(8.5), Inches(0.7),
             f"Stage {stage_num}: {stage_name}", size=32, color=TEXT_MAIN, bold=True)
    add_badge(s, Inches(10.6), Inches(1.08), status_text, status_color)

    col_y = Inches(1.95)
    col_h = Inches(4.6)
    col_w = Inches(5.95)
    left_x = Inches(0.55)
    right_x = Inches(6.85)

    add_rect(s, left_x, col_y, col_w, col_h, BG_PANEL, radius=0.05)
    add_text(s, left_x + Inches(0.4), col_y + Inches(0.35), col_w - Inches(0.8), Inches(0.5),
             "IN CODING", size=14, color=TEXT_FAINT, bold=True)
    add_text(s, left_x + Inches(0.4), col_y + Inches(0.8), col_w - Inches(0.8), Inches(0.6),
             coding_title, size=19, color=TEXT_DIM, bold=True, line_spacing=1.2)
    add_bullets(s, left_x + Inches(0.4), col_y + Inches(1.55), col_w - Inches(0.8), Inches(2.8),
                coding_bullets, size=14.5, color=TEXT_DIM, space_after=12)

    add_rect(s, right_x, col_y, col_w, col_h, BG_PANEL2, radius=0.05)
    add_rect(s, right_x, col_y, Inches(0.08), col_h, status_color)
    add_text(s, right_x + Inches(0.4), col_y + Inches(0.35), col_w - Inches(0.8), Inches(0.5),
             "ON MINDS, IN LAUNCHMINDS", size=14, color=status_color, bold=True)
    add_text(s, right_x + Inches(0.4), col_y + Inches(0.8), col_w - Inches(0.8), Inches(0.6),
             lm_title, size=19, color=TEXT_MAIN, bold=True, line_spacing=1.2)
    add_bullets(s, right_x + Inches(0.4), col_y + Inches(1.55), col_w - Inches(0.8), Inches(2.8),
                lm_bullets, size=14.5, color=TEXT_MAIN, space_after=12)

    add_text(s, Inches(0.55), Inches(6.75), Inches(11.8), Inches(0.4),
             footer_line, size=13.5, color=TEXT_FAINT, italic=True)
    add_footer(s)
    set_notes(s, notes)
    return s


# Slide 5 — Stage 1
stage_slide(
    "The Evolution · Stage 1", "04 / 07", 1, "Assisted Completion", "DONE", TEXT_FAINT,
    "Inline suggestions, not action",
    [
        "Suggestions only — human reviews every one",
        "Human stays in control of the outcome",
    ],
    "Minds-assisted campaign drafting",
    [
        "Minds drafts campaign structures, task mechanics, and copy",
        "Grounded in trusted project intelligence maintained by LaunchMinds",
        "Human approves every output",
    ],
    "Speed and consistency — but the human still owns every step.",
    "Stage one is assisted completion. A Mind drafts campaign structures, task mechanics, and copy grounded "
    "in the trusted project intelligence maintained by LaunchMinds. The operator reviews and approves every "
    "output. The value at this stage is speed and consistency, but the human still owns every step of the "
    "workflow.",
)

# Slide 6 — Stage 2
stage_slide(
    "The Evolution · Stage 2", "05 / 07", 2, "Single-Session Execution", "TODAY", ACCENT,
    "Full task, one sitting, human watching",
    [
        "One session completes an entire task",
        "A human observes, doesn't drive step by step",
    ],
    "One session, one Mind, full pipeline",
    [
        "Registration + planning + deployment — no handoffs between tools",
        "Runs through campaign-specific Skills, Tools, state, and approval gates on Minds",
        "Minds supplies the persistent runtime and execution; LaunchMinds supplies the campaign intelligence and permissions",
    ],
    "The Mind stops only suggesting and starts completing real operational work.",
    "Stage two is single-session execution. One persistent Mind can take a project from registration and "
    "briefing through campaign planning and deployment preparation without handoffs between different tools. "
    "LaunchMinds supplies the campaign-specific intelligence, Skills, state, and permissions. Minds supplies "
    "the persistent agent runtime and tool execution. This is where the Mind stops only suggesting and starts "
    "completing real operational work — while a human still supervises the session and approves high-impact "
    "actions.",
)

# ---------------------------------------------------------------------------
# Slide 7 — Stage 3: Persistent, Bounded Campaign Operations
# ---------------------------------------------------------------------------
s = new_slide()
add_kicker(s, "The Evolution · Stage 3", "06 / 07")
add_text(s, Inches(0.55), Inches(0.98), Inches(9.6), Inches(1.35),
         "Stage 3: Persistent, Bounded\nCampaign Operations", size=28, color=TEXT_MAIN, bold=True, line_spacing=1.15)
add_badge(s, Inches(10.35), Inches(1.08), "NEXT MILESTONE", ACCENT2)
add_bullets(s, Inches(0.55), Inches(2.55), Inches(11.8), Inches(3.6), [
    "Long-running project workspace with shared campaign state and evidence",
    "A dynamic team of Minds assembled around each objective",
    "Approval gates for launch, budget changes, and settlement",
    "Verify before reward; recover from failure; log every action",
    "Every outcome improves the next campaign",
], size=17.5, color=TEXT_MAIN, space_after=16)
add_text(s, Inches(0.55), Inches(6.35), Inches(11.8), Inches(0.5),
         "The shape of the team follows the work — not a fixed pipeline of roles.",
         size=14.5, color=ACCENT2, italic=True)
add_footer(s)
set_notes(s, "Stage three is not simply about adding more agents. It is the shift from single-session "
             "execution to persistent, bounded campaign operations. Each project gets a continuously "
             "maintained operating state: its objectives, constraints, budgets, active campaigns, participant "
             "evidence, pending approvals, unresolved work, and results. Minds provides the persistent agents, "
             "memory, tools, and coordination. LaunchMinds provides the campaign control plane: the trusted "
             "project intelligence, action permissions, approval gates, verification rules, and settlement "
             "logic. A dynamic team of Minds can form around each campaign — researching, planning, executing, "
             "monitoring, verifying, and recovering — without being locked into a fixed pipeline or permanent "
             "set of roles. The shape of the team follows the work. Humans define the mandate and approve "
             "high-impact actions. Routine operations continue autonomously within those boundaries. Every "
             "action is logged, participation is verified before rewards are released, and real outcomes "
             "update the next plan. That closes the campaign operations loop: plan, approve, execute, verify, "
             "settle, and learn.")

# ---------------------------------------------------------------------------
# Slide 8 — Closing / Vision
# ---------------------------------------------------------------------------
s = new_slide()
add_kicker(s, "Where This Goes", "07 / 07")
add_text(s, Inches(0.55), Inches(1.0), Inches(11.5), Inches(0.7),
         "Built on Minds — with Minds' community", size=32, color=TEXT_MAIN, bold=True)

timeline_y = Inches(2.0)
stages = [
    ("Stage 1", "Minds-assisted drafting", "DONE", TEXT_FAINT),
    ("Stage 2", "One Mind, full session", "TODAY", ACCENT),
    ("Stage 3", "Persistent, bounded ops", "NEXT", ACCENT2),
]
tx = [Inches(0.55), Inches(4.75), Inches(8.95)]
tw = Inches(3.65)
for (label, desc, status, c), x in zip(stages, tx):
    add_rect(s, x, timeline_y, tw, Inches(1.7), BG_PANEL, radius=0.08)
    add_rect(s, x, timeline_y, tw, Inches(0.07), c)
    add_text(s, x + Inches(0.3), timeline_y + Inches(0.25), tw - Inches(0.6), Inches(0.4),
             label, size=15, color=c, bold=True)
    add_text(s, x + Inches(0.3), timeline_y + Inches(0.65), tw - Inches(0.6), Inches(0.6),
             desc, size=17, color=TEXT_MAIN, bold=True, line_spacing=1.15)
    add_badge(s, x + Inches(0.3), timeline_y + Inches(1.25), status, c)

add_text(s, Inches(0.55), Inches(4.0), Inches(12.2), Inches(0.9),
         "Not an integration bolted on after the fact — built on Minds from day one.",
         size=19, color=ACCENT2, italic=True, line_spacing=1.3)
add_text(s, Inches(0.55), Inches(4.75), Inches(12.2), Inches(1.5),
         "Minds makes agents persistent and capable.\nLaunchMinds makes campaign operations accountable.",
         size=22, color=TEXT_MAIN, bold=True, line_spacing=1.25)
add_text(s, Inches(0.55), Inches(6.15), Inches(11.5), Inches(0.6),
         "We'd like to keep building that answer together — with your platform, and with your community.",
         size=15.5, color=TEXT_DIM, italic=True)
add_text(s, Inches(0.55), Inches(6.72), Inches(6), Inches(0.4),
         "Thank you.", size=16, color=TEXT_DIM, italic=True)
add_footer(s, "LaunchMinds · 2026")
set_notes(s, "So: two stages built and live today, one stage actively in motion, built on Minds from day one — "
             "not as an integration bolted on after the fact, but as the foundation. Minds makes agents "
             "persistent and capable. LaunchMinds makes campaign operations accountable. We're showing you "
             "this because we think it's a good answer to the question every platform has to answer "
             "eventually — what gets built on top of us? We'd like to keep building that answer together, "
             "with your platform and your community. Thanks for the time.")

prs.save("/tmp/claude-1000/-home-ubuntu/5e8a870b-0a20-4672-b81f-1aebc813a8e3/scratchpad/launchminds-deck/launchminds-deck.pptx")
print("Saved.")
